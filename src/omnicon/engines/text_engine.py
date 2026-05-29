# Copyright (c) 2026 Chenglin Qiu (SHC - Super Han Chinese). All rights reserved.
# Licensed under the OmniCon Proprietary License. See LICENSE for details.
"""Text extraction and data format conversion engine.

Handles text extraction from Office documents (DOCX, PPTX, XLSX) and
tabular data format conversions (XLSX <-> CSV) using python-docx,
python-pptx, openpyxl, and pandas.
"""

import logging
from pathlib import Path

from omnicon.core.job import ConversionJob
from omnicon.engines.base import BaseEngine, EngineError, UnsupportedRouteError

logger = logging.getLogger(__name__)


class TextEngine(BaseEngine):
    """Handles text extraction and data format conversions.

    Supports extracting plain text from DOCX, PPTX, and XLSX files,
    as well as converting between XLSX and CSV formats.
    """

    priority = 10

    SUPPORTED_ROUTES: set[tuple[str, str]] = {
        ("docx", "txt"),
        ("pptx", "txt"),
        ("xlsx", "txt"),
        ("xlsx", "csv"),
        ("csv", "xlsx"),
    }

    def can_convert(self, src_fmt: str, dst_fmt: str) -> bool:
        """Check if this engine supports the given conversion route."""
        return (src_fmt.lower(), dst_fmt.lower()) in self.SUPPORTED_ROUTES

    def convert(self, job: ConversionJob) -> Path:
        """Execute the conversion. Must be thread-safe.

        Args:
            job: The conversion job with input path, output format, and options.

        Returns:
            Path to the converted output file.

        Raises:
            EngineError: If conversion fails.
            UnsupportedRouteError: If the route is not supported by this engine.
        """
        route = (job.src_format, job.output_format.lower())
        match route:
            case ("docx", "txt"):
                return self._docx_to_text(job)
            case ("pptx", "txt"):
                return self._pptx_to_text(job)
            case ("xlsx", "txt"):
                return self._xlsx_to_text(job)
            case ("xlsx", "csv"):
                return self._xlsx_to_csv(job)
            case ("csv", "xlsx"):
                return self._csv_to_xlsx(job)
            case _:
                raise UnsupportedRouteError(
                    f"Route {route} not supported by {self.__class__.__name__}"
                )

    def _docx_to_text(self, job: ConversionJob) -> Path:
        """Extract text from a DOCX file using python-docx.

        Extracts paragraphs, tables, headers, and footers in document order.

        Args:
            job: The conversion job.

        Returns:
            Path to the output text file.

        Raises:
            EngineError: If text extraction fails.
        """
        output_path = job.expected_output_path
        logger.info("Extracting text from %s -> txt", job.input_path.name)

        try:
            from docx import Document
            from docx.table import Table
            from docx.text.paragraph import Paragraph

            doc = Document(str(job.input_path))
            parts: list[str] = []

            # Headers from first section
            for section in doc.sections:
                header_text = "\n".join(
                    p.text for p in section.header.paragraphs if p.text.strip()
                )
                if header_text:
                    parts.append(f"[Header]\n{header_text}")
                    break  # just first section header

            # Body: paragraphs and tables in document order
            for element in doc.element.body:
                tag = element.tag.split("}")[-1]  # strip namespace
                if tag == "p":
                    para = Paragraph(element, doc)
                    if para.text.strip():
                        parts.append(para.text)
                elif tag == "tbl":
                    table = Table(element, doc)
                    for row in table.rows:
                        row_text = "\t".join(cell.text for cell in row.cells)
                        parts.append(row_text)

            # Footers from first section
            for section in doc.sections:
                footer_text = "\n".join(
                    p.text for p in section.footer.paragraphs if p.text.strip()
                )
                if footer_text:
                    parts.append(f"\n[Footer]\n{footer_text}")
                    break

            output_path.write_text("\n\n".join(parts), encoding="utf-8")
        except Exception as exc:
            raise EngineError(f"DOCX text extraction failed: {exc}") from exc

        job.progress = 100
        return output_path

    def _pptx_to_text(self, job: ConversionJob) -> Path:
        """Extract text from a PPTX file using python-pptx.

        Extracts text from text frames, tables, and grouped shapes,
        plus speaker notes for each slide.

        Args:
            job: The conversion job.

        Returns:
            Path to the output text file.

        Raises:
            EngineError: If text extraction fails.
        """
        output_path = job.expected_output_path
        logger.info("Extracting text from %s -> txt", job.input_path.name)

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            def _extract_shapes(shapes: "Iterable", lines: list[str]) -> None:
                """Recursively extract text from shapes, tables, and groups."""
                for shape in shapes:
                    # Text frames (titles, body, text boxes)
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                lines.append(para.text)
                    # Tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = "\t".join(cell.text for cell in row.cells)
                            lines.append(row_text)
                    # Grouped shapes — recurse into child shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _extract_shapes(shape.shapes, lines)

            prs = Presentation(str(job.input_path))
            lines: list[str] = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"--- Slide {i + 1} ---")
                _extract_shapes(slide.shapes, lines)
                # Speaker notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        lines.append("")
                        lines.append("[Speaker Notes]")
                        lines.append(notes_text)
                lines.append("")  # blank line between slides
            output_path.write_text("\n".join(lines), encoding="utf-8")
        except Exception as exc:
            raise EngineError(f"PPTX text extraction failed: {exc}") from exc

        job.progress = 100
        return output_path

    def _xlsx_to_text(self, job: ConversionJob) -> Path:
        """Extract text from an XLSX file using openpyxl.

        Args:
            job: The conversion job.

        Returns:
            Path to the output text file.

        Raises:
            EngineError: If text extraction fails.
        """
        output_path = job.expected_output_path
        logger.info("Extracting text from %s -> txt", job.input_path.name)

        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(job.input_path), read_only=True, data_only=True)
            lines: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"--- {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    lines.append(
                        "\t".join(str(cell) if cell is not None else "" for cell in row)
                    )
            wb.close()
            output_path.write_text("\n".join(lines), encoding="utf-8")
        except Exception as exc:
            raise EngineError(f"XLSX text extraction failed: {exc}") from exc

        job.progress = 100
        return output_path

    def _xlsx_to_csv(self, job: ConversionJob) -> Path:
        """Convert an XLSX file to CSV using pandas.

        For multi-sheet workbooks, each sheet is exported as a separate CSV file
        named {stem}_{SheetName}.csv. Returns the path to the first CSV.
        Single-sheet workbooks produce a single {stem}.csv.

        Args:
            job: The conversion job.

        Returns:
            Path to the output CSV file (first sheet if multi-sheet).

        Raises:
            EngineError: If conversion fails.
        """
        output_path = job.expected_output_path
        logger.info("Converting %s -> csv", job.input_path.name)

        try:
            import pandas as pd

            sheets = pd.read_excel(job.input_path, engine="openpyxl", sheet_name=None)
            if len(sheets) == 1:
                df = next(iter(sheets.values()))
                df.to_csv(output_path, index=False, encoding="utf-8")
            else:
                first_path: Path | None = None
                for sheet_name, df in sheets.items():
                    safe_name = str(sheet_name).replace(" ", "_")
                    sheet_path = job.output_dir / f"{job.input_path.stem}_{safe_name}.csv"
                    df.to_csv(sheet_path, index=False, encoding="utf-8")
                    if first_path is None:
                        first_path = sheet_path
                output_path = first_path  # type: ignore[assignment]
        except Exception as exc:
            raise EngineError(f"XLSX to CSV conversion failed: {exc}") from exc

        job.progress = 100
        return output_path

    def _csv_to_xlsx(self, job: ConversionJob) -> Path:
        """Convert a CSV file to XLSX using pandas.

        Args:
            job: The conversion job.

        Returns:
            Path to the output XLSX file.

        Raises:
            EngineError: If conversion fails.
        """
        output_path = job.expected_output_path
        logger.info("Converting %s -> xlsx", job.input_path.name)

        try:
            import pandas as pd

            df = pd.read_csv(job.input_path, encoding="utf-8")
            df.to_excel(output_path, index=False, engine="openpyxl")
        except Exception as exc:
            raise EngineError(f"CSV to XLSX conversion failed: {exc}") from exc

        job.progress = 100
        return output_path
