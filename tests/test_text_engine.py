"""Tests for TextEngine conversions."""

from pathlib import Path

import pytest

from omnicon.core.job import ConversionJob
from omnicon.engines.base import EngineError, UnsupportedRouteError
from omnicon.engines.text_engine import TextEngine


@pytest.fixture
def engine() -> TextEngine:
    return TextEngine()


def make_job(input_path: Path, output_format: str, output_dir: Path) -> ConversionJob:
    return ConversionJob(
        input_path=input_path,
        output_format=output_format,
        output_dir=output_dir,
    )


# ---------------------------------------------------------------------------
# Fixture generators — create test files programmatically
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    """Create a minimal DOCX with paragraphs and a table using python-docx."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello OmniCon!")
    doc.add_paragraph("Second paragraph for testing extraction.")

    # Add a table with header row + data rows
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Product"
    table.cell(0, 1).text = "Price"
    table.cell(1, 0).text = "Widget"
    table.cell(1, 1).text = "9.99"
    table.cell(2, 0).text = "Gadget"
    table.cell(2, 1).text = "19.99"

    doc.add_paragraph("Paragraph after table.")
    path = tmp_path / "sample.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a PPTX with two slides, speaker notes, and a table."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = f"Slide {i + 1} Title"
        slide.placeholders[1].text = f"Content for slide {i + 1}"

    # Add speaker notes to the first slide
    notes_slide = prs.slides[0].notes_slide
    notes_slide.notes_text_frame.text = "Remember to greet the audience"

    # Add a table to the second slide
    slide2 = prs.slides[1]
    table_shape = slide2.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(3), width=Inches(4), height=Inches(1.5),
    )
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Score"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "95"

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    """Create a minimal XLSX with data using openpyxl."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", 30, "Perth"])
    ws.append(["Bob", 25, "Sydney"])
    path = tmp_path / "sample.xlsx"
    wb.save(str(path))
    wb.close()
    return path


@pytest.fixture
def sample_csv(tmp_path: Path) -> Path:
    """Create a minimal CSV file."""
    path = tmp_path / "sample.csv"
    path.write_text("Name,Age,City\nAlice,30,Perth\nBob,25,Sydney\n", encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# can_convert tests
# ---------------------------------------------------------------------------


def test_can_convert_supported(engine: TextEngine) -> None:
    assert engine.can_convert("docx", "txt") is True
    assert engine.can_convert("pptx", "txt") is True
    assert engine.can_convert("xlsx", "txt") is True
    assert engine.can_convert("xlsx", "csv") is True
    assert engine.can_convert("csv", "xlsx") is True


def test_can_convert_unsupported(engine: TextEngine) -> None:
    assert engine.can_convert("pdf", "txt") is False
    assert engine.can_convert("docx", "pdf") is False
    assert engine.can_convert("txt", "docx") is False
    assert engine.can_convert("csv", "txt") is False


def test_can_convert_case_insensitive(engine: TextEngine) -> None:
    assert engine.can_convert("DOCX", "TXT") is True
    assert engine.can_convert("Xlsx", "Csv") is True
    assert engine.can_convert("CSV", "XLSX") is True


# ---------------------------------------------------------------------------
# DOCX -> TXT
# ---------------------------------------------------------------------------


def test_docx_to_text(engine: TextEngine, sample_docx: Path, tmp_path: Path) -> None:
    """DOCX to TXT extraction produces a valid text file with paragraphs and tables."""
    job = make_job(sample_docx, "txt", tmp_path)
    result = engine.convert(job)

    assert result.exists()
    assert result.suffix == ".txt"
    content = result.read_text(encoding="utf-8")
    assert "Hello OmniCon!" in content
    assert "Second paragraph" in content
    # Table content should be extracted
    assert "Product" in content
    assert "Price" in content
    assert "Widget" in content
    assert "9.99" in content
    assert "Gadget" in content
    assert "19.99" in content
    # Paragraph after the table
    assert "Paragraph after table." in content
    assert job.progress == 100


def test_docx_to_text_output_filename(
    engine: TextEngine, sample_docx: Path, tmp_path: Path
) -> None:
    """Output filename follows the stem.output_format convention."""
    job = make_job(sample_docx, "txt", tmp_path)
    result = engine.convert(job)
    assert result.name == "sample.txt"


# ---------------------------------------------------------------------------
# PPTX -> TXT
# ---------------------------------------------------------------------------


def test_pptx_to_text(engine: TextEngine, sample_pptx: Path, tmp_path: Path) -> None:
    """PPTX to TXT extraction produces text with slide markers, notes, and tables."""
    job = make_job(sample_pptx, "txt", tmp_path)
    result = engine.convert(job)

    assert result.exists()
    assert result.suffix == ".txt"
    content = result.read_text(encoding="utf-8")
    # Slide markers and basic text
    assert "--- Slide 1 ---" in content
    assert "--- Slide 2 ---" in content
    assert "Slide 1 Title" in content
    assert "Content for slide 2" in content
    # Speaker notes from slide 1
    assert "[Speaker Notes]" in content
    assert "Remember to greet the audience" in content
    # Table data from slide 2
    assert "Name\tScore" in content
    assert "Alice\t95" in content
    assert job.progress == 100


# ---------------------------------------------------------------------------
# XLSX -> TXT
# ---------------------------------------------------------------------------


def test_xlsx_to_text(engine: TextEngine, sample_xlsx: Path, tmp_path: Path) -> None:
    """XLSX to TXT extraction produces a tab-delimited text file."""
    job = make_job(sample_xlsx, "txt", tmp_path)
    result = engine.convert(job)

    assert result.exists()
    assert result.suffix == ".txt"
    content = result.read_text(encoding="utf-8")
    assert "--- Sheet1 ---" in content
    assert "Alice" in content
    assert "Perth" in content
    assert job.progress == 100


# ---------------------------------------------------------------------------
# XLSX -> CSV
# ---------------------------------------------------------------------------


def test_xlsx_to_csv(engine: TextEngine, sample_xlsx: Path, tmp_path: Path) -> None:
    """XLSX to CSV conversion produces a valid CSV file."""
    job = make_job(sample_xlsx, "csv", tmp_path)
    result = engine.convert(job)

    assert result.exists()
    assert result.suffix == ".csv"
    content = result.read_text(encoding="utf-8")
    assert "Name" in content
    assert "Alice" in content
    assert "Bob" in content
    assert job.progress == 100


def test_xlsx_to_csv_content_integrity(
    engine: TextEngine, sample_xlsx: Path, tmp_path: Path
) -> None:
    """CSV output preserves all rows from the XLSX source."""
    job = make_job(sample_xlsx, "csv", tmp_path)
    result = engine.convert(job)
    lines = result.read_text(encoding="utf-8").strip().splitlines()
    # Header + 2 data rows
    assert len(lines) == 3


# ---------------------------------------------------------------------------
# CSV -> XLSX
# ---------------------------------------------------------------------------


def test_csv_to_xlsx(engine: TextEngine, sample_csv: Path, tmp_path: Path) -> None:
    """CSV to XLSX conversion produces a valid XLSX file."""
    job = make_job(sample_csv, "xlsx", tmp_path)
    result = engine.convert(job)

    assert result.exists()
    assert result.suffix == ".xlsx"
    assert result.stat().st_size > 0
    assert job.progress == 100


def test_csv_to_xlsx_content_integrity(
    engine: TextEngine, sample_csv: Path, tmp_path: Path
) -> None:
    """XLSX output preserves all rows from the CSV source."""
    import openpyxl

    job = make_job(sample_csv, "xlsx", tmp_path)
    result = engine.convert(job)

    wb = openpyxl.load_workbook(str(result), read_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    wb.close()

    assert len(rows) == 3  # Header + 2 data rows
    assert rows[0] == ("Name", "Age", "City")
    assert rows[1][0] == "Alice"


# ---------------------------------------------------------------------------
# Roundtrip tests
# ---------------------------------------------------------------------------


def test_xlsx_csv_roundtrip(engine: TextEngine, sample_xlsx: Path, tmp_path: Path) -> None:
    """XLSX -> CSV -> XLSX roundtrip preserves data."""
    import openpyxl

    # XLSX -> CSV
    csv_dir = tmp_path / "csv_out"
    csv_dir.mkdir()
    job1 = make_job(sample_xlsx, "csv", csv_dir)
    csv_path = engine.convert(job1)

    # CSV -> XLSX
    xlsx_dir = tmp_path / "xlsx_out"
    xlsx_dir.mkdir()
    job2 = make_job(csv_path, "xlsx", xlsx_dir)
    xlsx_path = engine.convert(job2)

    wb = openpyxl.load_workbook(str(xlsx_path), read_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    wb.close()

    assert len(rows) == 3
    assert rows[0] == ("Name", "Age", "City")


# ---------------------------------------------------------------------------
# Error handling
# ---------------------------------------------------------------------------


def test_unsupported_route_raises(engine: TextEngine, sample_docx: Path, tmp_path: Path) -> None:
    """Requesting an unsupported route raises UnsupportedRouteError."""
    job = make_job(sample_docx, "pdf", tmp_path)
    with pytest.raises(UnsupportedRouteError):
        engine.convert(job)


def test_invalid_docx_raises_engine_error(engine: TextEngine, tmp_path: Path) -> None:
    """Corrupt DOCX file raises EngineError."""
    bad_file = tmp_path / "corrupt.docx"
    bad_file.write_bytes(b"this is not a valid docx")
    job = make_job(bad_file, "txt", tmp_path)
    with pytest.raises(EngineError):
        engine.convert(job)


def test_invalid_xlsx_raises_engine_error(engine: TextEngine, tmp_path: Path) -> None:
    """Corrupt XLSX file raises EngineError."""
    bad_file = tmp_path / "corrupt.xlsx"
    bad_file.write_bytes(b"this is not a valid xlsx")
    job = make_job(bad_file, "txt", tmp_path)
    with pytest.raises(EngineError):
        engine.convert(job)


# ---------------------------------------------------------------------------
# Priority
# ---------------------------------------------------------------------------


def test_priority(engine: TextEngine) -> None:
    assert engine.priority == 10
